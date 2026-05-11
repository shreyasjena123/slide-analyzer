from dataclasses import dataclass, field
from pathlib import Path
import fitz  # PyMuPDF


@dataclass
class SlideData:
    num: int
    title: str
    body_text: str
    notes: str
    image_count: int
    image_positions: list[tuple]  # (x0, y0, x1, y1) normalized 0-1
    text_box_positions: list[tuple]  # (x0, y0, x1, y1) normalized 0-1, paired with text
    slide_width: float
    slide_height: float

    def has_text_image_separation(self) -> bool:
        """True if images and text blocks are on opposite halves of the slide."""
        if not self.image_positions or not self.text_box_positions:
            return False
        img_centers_x = [((x0 + x1) / 2) for x0, y0, x1, y1 in self.image_positions]
        txt_centers_x = [((x0 + x1) / 2) for x0, y0, x1, y1 in self.text_box_positions]
        avg_img_x = sum(img_centers_x) / len(img_centers_x)
        avg_txt_x = sum(txt_centers_x) / len(txt_centers_x)
        # Flag if average centers are more than 40% of slide width apart
        return abs(avg_img_x - avg_txt_x) > 0.4

    _TRANSITION_KEYWORDS = {
        "overview", "agenda", "outline", "objectives", "objective",
        "questions", "thank", "summary", "contents", "today", "recap",
        "review", "introduction", "intro", "break", "discussion",
    }

    def slide_type(self) -> str:
        """Heuristic classification: title/transition, section, or content."""
        title_words = set(self.title.lower().split())
        if title_words & self._TRANSITION_KEYWORDS:
            return "section/transition"
        if len(self.body_text.strip()) < 60 and self.image_count == 0:
            return "title/transition"
        return "content"

    def to_prompt_text(self) -> str:
        lines = [
            f"Slide {self.num}",
            f"Slide type: {self.slide_type()} (title/transition slides are intentionally sparse — be lenient)",
            f"Title: {self.title or '(no title)'}",
            f"Body text: {self.body_text or '(none)'}",
            f"Speaker notes: {self.notes or '(none)'}",
            f"Images on slide: {self.image_count}",
        ]
        if self.image_count > 0:
            lines.append(f"Image positions (normalized 0-1): {self.image_positions}")
            lines.append(f"Text box positions (normalized 0-1): {self.text_box_positions}")
            lines.append(f"Text and images on opposite halves: {self.has_text_image_separation()}")
        return "\n".join(lines)


def extract_pptx(path: str) -> list[SlideData]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    w = prs.slide_width
    h = prs.slide_height
    slides = []

    for i, slide in enumerate(prs.slides):
        title = ""
        body_parts = []
        notes = ""
        image_positions = []
        text_box_positions = []

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        for shape in slide.shapes:
            left = shape.left / w if w else 0
            top = shape.top / h if h else 0
            right = (shape.left + shape.width) / w if w else 0
            bottom = (shape.top + shape.height) / h if h else 0
            bbox = (round(left, 3), round(top, 3), round(right, 3), round(bottom, 3))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_positions.append(bbox)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape.name.lower().startswith("title") or (hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0):
                    title = text
                else:
                    body_parts.append(text)
                text_box_positions.append(bbox)

        slides.append(SlideData(
            num=i + 1,
            title=title,
            body_text="\n".join(body_parts),
            notes=notes,
            image_count=len(image_positions),
            image_positions=image_positions,
            text_box_positions=text_box_positions,
            slide_width=float(w),
            slide_height=float(h),
        ))

    return slides


def extract_pdf(path: str) -> list[SlideData]:
    doc = fitz.open(path)
    slides = []

    for i, page in enumerate(doc):
        w, h = page.rect.width, page.rect.height
        blocks = page.get_text("dict")["blocks"]

        title = ""
        body_parts = []
        text_box_positions = []
        image_positions = []

        for block in blocks:
            if block["type"] == 0:  # text
                text = " ".join(
                    span["text"] for line in block["lines"] for span in line["spans"]
                ).strip()
                if not text:
                    continue
                x0, y0, x1, y1 = block["bbox"]
                bbox = (round(x0/w, 3), round(y0/h, 3), round(x1/w, 3), round(y1/h, 3))
                # Heuristic: first large text block near top is title
                avg_size = sum(
                    span["size"] for line in block["lines"] for span in line["spans"]
                ) / max(sum(len(line["spans"]) for line in block["lines"]), 1)
                if not title and y0 < h * 0.25 and avg_size > 14:
                    title = text
                else:
                    body_parts.append(text)
                text_box_positions.append(bbox)
            elif block["type"] == 1:  # image
                x0, y0, x1, y1 = block["bbox"]
                image_positions.append((round(x0/w, 3), round(y0/h, 3), round(x1/w, 3), round(y1/h, 3)))

        slides.append(SlideData(
            num=i + 1,
            title=title,
            body_text="\n".join(body_parts),
            notes="",
            image_count=len(image_positions),
            image_positions=image_positions,
            text_box_positions=text_box_positions,
            slide_width=w,
            slide_height=h,
        ))

    return slides


def extract(path: str) -> list[SlideData]:
    p = Path(path)
    if p.suffix.lower() == ".pptx":
        return extract_pptx(path)
    elif p.suffix.lower() == ".pdf":
        return extract_pdf(path)
    else:
        raise ValueError(f"Unsupported file type: {p.suffix}")
